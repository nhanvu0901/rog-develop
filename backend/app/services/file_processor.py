from langchain_community.document_loaders import (
    Docx2txtLoader,
    UnstructuredExcelLoader,
    TextLoader,
    CSVLoader,
    PyPDFLoader,
)
import os
from pptx import Presentation

async def process_file(file_path: str, file_extension: str) -> str:
    """
    Process different types of files and extract text using LangChain
    """
    try:
        if file_extension.lower() == "docx":
            loader = Docx2txtLoader(file_path)
        elif file_extension.lower() in ["xlsx"]:
            loader = UnstructuredExcelLoader(file_path, mode="elements")
        elif file_extension.lower() in ["pptx"]:
            return extract_text_from_pptx(file_path).strip()
        elif file_extension.lower() in ["pdf"]:
            loader = PyPDFLoader(file_path)
        elif file_extension.lower() in ["csv"]:
            loader = CSVLoader(file_path)
        else:
            loader = TextLoader(file_path)
        
        documents = loader.load()
        extracted_text = ""
        
        for doc in documents:
            extracted_text += doc.page_content + "\n"
        
        return extracted_text.strip()
        
    except Exception as e:
        raise Exception(f"Error processing file: {str(e)}")
    # finally:
        # Clean up the uploaded file
        # if os.path.exists(file_path):
        #    os.remove(file_path) 
        
        
def extract_text_from_pptx(pptx_path):
    """
    Extract all text from a PowerPoint file
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        dict: Dictionary with slide numbers as keys and text content as values
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"File not found: {pptx_path}")
        
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Dictionary to store slide content
    slides_text = {}
    
    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        text_content = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_content.append(shape.text)
                
        # Extract text from tables
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text)
                    if row_text:
                        text_content.append(" | ".join(row_text))
        
        # Join all text elements with newlines
        slides_text[f"Slide {i}"] = "\n".join(text_content)
    
    return "\n".join(slides_text.values())